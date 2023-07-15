from flask import Flask, request, jsonify
from flask_sqlalchemy import SQLAlchemy
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from flask_cors import CORS

app = Flask(__name__)
CORS(app, origins='*')
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///database.db'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB maximum file size
db = SQLAlchemy(app)

# Database model
class PitchDeck(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    title = db.Column(db.String(255))
    content = db.Column(db.Text)

    def __repr__(self):
        return '<PitchDeck ' + self.title + '>'

def parse_pdf_content(pdf):
    text = []
    for page in pdf.pages:
        text.append(page.extract_text().strip())
    return '\n\n'.join(text)

def parse_ppt_content(presentation):
    slide_data = []
    for slide in presentation.slides:
        content = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                content.append(shape.text)
        slide_data.append(content)
    return '\n\n'.join(['\n'.join(content) for content in slide_data])

def parse_docx_content(doc):
    paragraphs = doc.paragraphs
    return '\n\n'.join([p.text.strip() for p in paragraphs])

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return 'No file part in the request', 400

    file = request.files['file']

    if file.content_length > app.config['MAX_CONTENT_LENGTH']:
        return 'File size exceeds the allowed limit', 400

    if file.mimetype == 'application/pdf':
        # Parse PDF file
        pdf = PdfReader(file)
        content = parse_pdf_content(pdf)

    elif file.mimetype in ['application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
        # Parse PowerPoint file
        prs = Presentation(file)
        content = parse_ppt_content(prs)

    elif file.mimetype in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
        # Parse Word document file
        doc = Document(file)
        content = parse_docx_content(doc)

    else:
        return 'Invalid file format. Please upload a PDF, PowerPoint, or Word document.', 400

    # Save to the database
    pitch_deck = PitchDeck(title=file.filename, content=content)
    db.session.add(pitch_deck)
    db.session.commit()

    return jsonify({
        'id': pitch_deck.id,
        'title': pitch_deck.title,
        'content': pitch_deck.content
    }), 200, {'Access-Control-Allow-Origin': '*'}

@app.route('/pitchdecks', methods=['GET'])
def get_pitch_decks():
    pitch_decks = PitchDeck.query.all()
    results = []
    for pitch_deck in pitch_decks:
        results.append({
            'id': pitch_deck.id,
            'title': pitch_deck.title,
            'content': pitch_deck.content
        })
    return jsonify(results), 200, {'Access-Control-Allow-Origin': '*'}

if __name__ == '__main__':
    db.create_all()
    app.run()
